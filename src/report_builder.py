import os
import pandas as pd
from sqlalchemy import create_engine

DB_URL = "postgresql://_shaidolla@localhost/process_lens"
engine = create_engine(DB_URL)

os.makedirs("outputs/charts", exist_ok=True)

print("Loading source data from previous phases...")
mann_whitney = pd.read_csv("outputs/mann_whitney_results.csv")
sensitivity = pd.read_csv("outputs/reports/financial_sensitivity.csv")
simulation = pd.read_csv("outputs/reports/simulation_results.csv")
feature_importance = pd.read_csv(
    "outputs/reports/classifier_feature_importance.csv", index_col=0
)

with open("sql/queries/05_cost_to_serve.sql") as f:
    cost_to_serve = pd.read_sql(f.read(), engine).head(5)

print(f"  Mann-Whitney results: {len(mann_whitney)} stages")
print(f"  Sensitivity scenarios: {len(sensitivity)}")
print(f"  Simulation scenarios: {len(simulation)}")
print(f"  Top cost-to-serve routes: {len(cost_to_serve)}")

import matplotlib.pyplot as plt

print("\nBuilding charts...")

stage_order = ["Order Approved", "Picked Packed and Shipped", "Delivered"]
medians = mann_whitney.set_index("activity").loc[stage_order, "median_on_time_d"]
fig, ax = plt.subplots(figsize=(6, 4))
ax.bar(stage_order, medians, color=["#4C72B0", "#4C72B0", "#C44E52"])
ax.set_ylabel("Median days (on-time orders)")
ax.set_title("Cycle Time by Stage")
plt.xticks(rotation=15)
plt.tight_layout()
plt.savefig("outputs/charts/stage_cycle_time.png", dpi=150)
plt.close()

labels = cost_to_serve["seller_state"] + " -> " + cost_to_serve["customer_state"]
fig, ax = plt.subplots(figsize=(6, 4))
ax.barh(labels, cost_to_serve["total_cost_to_serve"], color="#55A868")
ax.set_xlabel("Total cost to serve ($)")
ax.set_title("Top 5 Routes by Cost to Serve")
ax.invert_yaxis()
plt.tight_layout()
plt.savefig("outputs/charts/cost_to_serve.png", dpi=150)
plt.close()

fig, ax = plt.subplots(figsize=(5, 4))
ax.bar(simulation["scenario"], simulation["otd_rate"] * 100, color=["#C44E52", "#55A868"])
ax.set_ylabel("OTD rate (%)")
ax.set_title("Simulated OTD: As-Is vs To-Be")
ax.set_ylim(80, 100)
plt.tight_layout()
plt.savefig("outputs/charts/otd_comparison.png", dpi=150)
plt.close()

top_features = feature_importance.iloc[:, 0].head(5)
fig, ax = plt.subplots(figsize=(6, 4))
ax.barh(top_features.index, top_features.values, color="#8172B2")
ax.set_xlabel("Mean absolute SHAP value")
ax.set_title("Top 5 Predictive Features")
ax.invert_yaxis()
plt.tight_layout()
plt.savefig("outputs/charts/feature_importance.png", dpi=150)
plt.close()

print("  Charts saved to outputs/charts/")

from pptx import Presentation
from pptx.util import Inches, Pt

print("\nAssembling slide deck...")
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(4), Inches(12), Inches(1))
    tf2 = sub_box.text_frame
    tf2.text = subtitle
    tf2.paragraphs[0].font.size = Pt(20)


def add_content_slide(title, bullets, chart_path=None):
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    body_width = Inches(6) if chart_path else Inches(12)
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), body_width, Inches(5.5))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {bullet}"
        p.font.size = Pt(16)

    if chart_path:
        slide.shapes.add_picture(chart_path, Inches(6.8), Inches(1.3), width=Inches(6))


add_title_slide(
    "Process Lens — Order-to-Cash Business Case",
    "Where is the O2C fulfillment process breaking down, what does it cost, and what's the fix?",
)

add_content_slide(
    "Situation",
    [
        "96,455 real Olist e-commerce orders analyzed end-to-end",
        "Overall on-time delivery (OTD) rate: 91.9%",
        "Late order rate: 8.1% (7,826 orders)",
        "Carrier delivery is 3.9x longer than pick/pack -- 7.09 vs 1.84 days median",
    ],
    chart_path="outputs/charts/stage_cycle_time.png",
)

top_route_cost = cost_to_serve.iloc[0]["total_cost_to_serve"]
add_content_slide(
    "Complication",
    [
        "Delivered stage confirmed as the primary bottleneck (Mann-Whitney U, p<0.000001, large effect)",
        "Late orders take about 17 more days at this single stage than on-time orders",
        "About 18 seller-customer routes account for 80% of all late deliveries",
        f"Highest-cost route alone costs ${top_route_cost:,.0f}/year to serve",
    ],
    chart_path="outputs/charts/cost_to_serve.png",
)

base_row = sensitivity[sensitivity["scenario"].str.contains("Base")].iloc[0]
pess_row = sensitivity[sensitivity["scenario"].str.contains("Pessimistic")].iloc[0]
add_content_slide(
    "Answer",
    [
        "Recommendation: target carrier renegotiation on the top ~18 highest-cost routes",
        f"Base case: 3-year NPV ${base_row['npv']:,.0f}, payback in {base_row['payback_years']:.1f} years",
        f"Even the pessimistic scenario is NPV-positive (${pess_row['npv']:,.0f})",
        "Validated via discrete-event simulation before recommending, not just modeled on paper",
    ],
)

as_is_otd = simulation[simulation["scenario"] == "as_is"]["otd_rate"].values[0] * 100
to_be_otd = simulation[simulation["scenario"] == "to_be"]["otd_rate"].values[0] * 100
add_content_slide(
    "Validation",
    [
        f"Simulation calibrated to real data: {as_is_otd:.1f}% simulated vs 91.9% observed",
        f"To-be redesign: OTD improves from {as_is_otd:.1f}% to {to_be_otd:.1f}% ({to_be_otd - as_is_otd:.1f}pp)",
        "Confirmed via non-overlapping 95% confidence intervals across 20 replications",
        "Predictive classifier (AUC-ROC 0.74) flags at-risk orders at creation -- top 20% flagged catches 41.5% of late orders",
    ],
    chart_path="outputs/charts/otd_comparison.png",
)

prs.save("outputs/reports/process_lens_business_case.pptx")
print("\nDeck saved: outputs/reports/process_lens_business_case.pptx")